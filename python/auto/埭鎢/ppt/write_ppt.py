# coding:utf-8

'''
生成ppt樣式: p = Presentation()
選擇ppt樣式: slide = p.slide_layout[num] -num: 0-8
插入樣式到ppt的一頁中: one_ppt_slide = p.slides.add_slide(layout)
獲取標題 title = one_ppt_slide.placeholder[num]
獲取段落 content = one_ppt_slide.slide_layout[num]

段落共有三種方案:
    1. 直接寫內容: content.text = 內容
    2. 想讓每行都有不同的樣式: paragraph = content.text_frame.add_paragraph()
    3. 自定義一個文本框: box = slide.shapes.add_textbox(left, top, width, height)

表格:
    table = slide.shapes.add_table(cols, rows, left, top, width, height).table -> 獲取table對象
    table.cell(cols, rows).text = 內容

圖片:
    image = slide.shapes_add_picture(image_file圖片地址, left, top, width, height)
樣式:
    font.bold -文字加粗 = True
    font.italic -斜體 = True
    font.size -大小 -> Pt() => from pptx.util import Pt
    alignment -段落位置 -> PP_PARAGRAPH_ALIGNMENT => from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
    color.rgb -文字顏色 -> RGBColor => from pptx.dml.color import RGBColor
    font.underline -文字下划線 = True
文本框樣式: Inches() -> from pptx.util import Inches
'''
import pptx
from pptx.util import Pt, Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

p = pptx.Presentation()
layout = p.slide_layouts[1]  # 獲取段落: placeholder = slide.placeholders[1] -因為使用布局1, 有title, content, 所以用1獲取content
slide = p.slides.add_slide(layout)  # insert layout into one ppt slide

# placeholders[num] -num:0 title,  num:1 content

placeholder = slide.placeholders[1]  # 內容

# 段落添加內容: placeholder.text = 內容
#placeholder.text = '欢迎学习ppt制作\n欢迎学习python'

title = slide.placeholders[0]  # 標題
title.text = '题目'

# 在段落中定義多個段落
# 想讓段落中每一行都有各自的樣式 -初始化段落中 定義多個段落 -> paragraph = placeholder.text_frame.add_paragraph()
'''
text -內容
font.bold -文字加粗
font.italic -斜體
font.size -大小
alignment -段落位置
color.rgb -文字顏色
font.underline -文字下划線
'''
# 第一個段落
paragraph1 = placeholder.text_frame.add_paragraph()
paragraph1.text = '欢迎学习ppt制作'
paragraph1.bold = True
paragraph1.font.italic = True
paragraph1.font.size = Pt(16)  # from pptx.util import Pt
paragraph1.font.underline = True
paragraph1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  # from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

# 第二個段落
paragraph2 = placeholder.text_frame.add_paragraph()
paragraph2.text = '欢迎学习python'
paragraph2.font.size = Pt(32)
paragraph2.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

# 加入新的ppt分頁
layout = p.slide_layouts[6]  # title
slide = p.slides.add_slide(layout)

# 自定義文本框:
# box = slide.shapes.add_textbox(left, top, width, height) -left, top決定文本框位置, width, height決定文本框本身
left = top = width = height = Inches(5)  # 大小: Inches() -from pptx.util import Inches
box = slide.shapes.add_textbox(left, top, width, height)

# 生成大段落中的小分段
para = box.text_frame.add_paragraph()

para.text = 'this is a para test'  # 寫入內容
para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  # from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
para.font.size = Pt(32)  # from pptx.util import Pt
para.font.color.rgb = RGBColor(255, 255, 0)  # from pptx.dml.color import RGBColor
para.font.name = '微软雅黑'  # 具體字體樣式名稱

# 創建第三頁ppt
layout = p.slide_layouts[1]  # 定義布局
slide = p.slides.add_slide(layout)  # insert layout into one ppt_slide

rows = 10  # 10列
cols = 2  # 2行

# 表格位置
left = top = Inches(2)  # from pptx.util import Inches

# 表格大小
width = Inches(6.0)
height = Inches(1.0)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table  # 獲取table對象

# 通過循環, 添加索引
for index, _ in enumerate(range(rows)):  # 列
    for sub_index in range(cols):  # 行
        # index, sub_index 對應表格中具體的位置
        table.cell(index, sub_index).text = '%s:%s' % (index, sub_index)  # 插入內容.text = 內容

layout = p.slide_layouts[6]
slide = p.slides.add_slide(layout)

image = slide.shapes.add_picture(
    image_file='logo2020.png',
    left=Inches(1),
    top=Inches(1),
    width=Inches(6),
    height=Inches(4)
)

p.save('test2.ppt')
