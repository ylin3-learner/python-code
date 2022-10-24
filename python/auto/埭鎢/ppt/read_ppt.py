# coding:utf-8

# 讀取ppt(不包含圖片)
'''
獲取整體的ppt對象: p = Presentation(ppt_address)
                p.slides -> 返回一個幻燈列表, 獲取所有ppt內容
                for i in slide.shapes:　返回具體的ppt形狀
獲取文本內容:
        shape.has_text_frame() -> 判斷是否是文本類型
        shape.text_frame.text() -> 獲取文本內容
獲取表格內容:
        shape.has_table -> 判斷是否是表格類型
        shape.table.iter_cells() -> 獲取文本內容列表
'''
import pptx

p = pptx.Presentation('test2.ppt')
for slide in p.slides:  # 獲取ppt_slides
    print(slide)
    for shape in slide.shapes:
        print(shape)
        if shape.has_text_frame:  # 判斷是否是文本類型
            print(shape.text_frame.text)  # 打印文本內容
        if shape.has_table:  # 判斷是否是表格類型
            for cell in shape.table.iter_cells():
                print(cell.text)  # 打印表格中文本內容

