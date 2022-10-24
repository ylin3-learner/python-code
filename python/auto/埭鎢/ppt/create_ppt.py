# coding:utf-8

# ppt的創建:
'''
pip install python -pptx
import pptx
'''
# 創建空ppt對象: p = Presentation()
# 布局: layout = p.slide_layouts(num) -num: 0-8九種布局選項
# 生成一頁ppt: slide = p.slides.add_slides(layout)  -slide 為一頁ppt的對象
# 保存ppt具體的地址: p.save(pptx)


# 生成一頁空ppt
import pptx  # pip install python-pptx

p = pptx.Presentation()  # 生成ppt对象

layout = p.slide_layouts[8]  # 樣式

# 0 一個title, 一個sub_title
# 1 有一個title, 一個content
# 7 one title, one sub_title, one content
slide = p.slides.add_slide(layout)  # 把樣式插入到ppt中生成第一頁


# 在一頁ppt中添加內容

p.save('test1.ppt')
